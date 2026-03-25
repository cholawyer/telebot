"""
책 사진 → 4세 수준 요약 PPTX 봇
텔레그램으로 책 사진 보내면 → PPTX 생성 후 전송
"""

import os
import json
import time
import base64
import requests
import tempfile
import urllib.parse
from io import BytesIO
from pathlib import Path

import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

BOT_TOKEN = "8729968345:AAEcJcUiaKvsa0aOuOmpSaafOwcgRB3P1Oo"
API_URL = f"https://api.telegram.org/bot{BOT_TOKEN}"
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")

client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# ── 텔레그램 API ──────────────────────────────────────────────

def tg_get(method, params=None):
    r = requests.get(f"{API_URL}/{method}", params=params, timeout=30)
    return r.json()

def tg_post(method, data=None, files=None):
    r = requests.post(f"{API_URL}/{method}", data=data, files=files, timeout=60)
    return r.json()

def send_message(chat_id, text):
    tg_post("sendMessage", {"chat_id": chat_id, "text": text})

def send_document(chat_id, filepath, caption=""):
    with open(filepath, "rb") as f:
        tg_post("sendDocument", {"chat_id": chat_id, "caption": caption},
                files={"document": f})

def download_photo(file_id):
    """텔레그램 파일 다운로드 → bytes 반환"""
    info = tg_get("getFile", {"file_id": file_id})
    file_path = info["result"]["file_path"]
    url = f"https://api.telegram.org/file/bot{BOT_TOKEN}/{file_path}"
    r = requests.get(url, timeout=30)
    return r.content

# ── Claude: 책 분석 ──────────────────────────────────────────

def analyze_book_image(image_bytes):
    """책 사진 → 6~8개 슬라이드 내용 생성"""
    b64 = base64.standard_b64encode(image_bytes).decode("utf-8")

    prompt = """이 책의 내용을 분석해줘.

책의 텍스트를 읽고 4살 아이 수준으로 이해하기 쉽게 요약해서 슬라이드 6~8장을 만들어줘.

각 슬라이드는:
- 핵심 내용을 짧고 재미있게 (20자 이내)
- 관련 이미지 검색 키워드 (영어로, 2~3단어, 어린이 삽화 스타일)

반드시 아래 JSON 형식으로만 답해줘 (다른 텍스트 없이):
{
  "title": "책 제목 또는 주제",
  "slides": [
    {
      "caption": "4살 수준 짧은 설명",
      "keyword": "english keyword for image search"
    }
  ]
}"""

    msg = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=1500,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/jpeg",
                        "data": b64
                    }
                },
                {"type": "text", "text": prompt}
            ]
        }]
    )

    raw = msg.content[0].text.strip()
    # JSON 블록 추출
    if "```" in raw:
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())

# ── 이미지 검색 ───────────────────────────────────────────────

def search_image(keyword):
    """키워드로 이미지 검색 후 bytes 반환 (LoremFlickr 사용)"""
    # 첫 번째 단어만 사용 (더 정확한 결과)
    first_word = keyword.split()[0] if keyword else "nature"

    try:
        url = f"https://loremflickr.com/800/600/{urllib.parse.quote(first_word)}"
        r = requests.get(url, timeout=20, allow_redirects=True)
        if r.status_code == 200 and r.content[:3] in (b'\xff\xd8\xff', b'\x89PN', b'GIF'):
            return r.content
    except Exception:
        pass

    # 폴백: 전체 키워드
    try:
        full = urllib.parse.quote(keyword.replace(" ", ","))
        r = requests.get(f"https://loremflickr.com/800/600/{full}", timeout=15, allow_redirects=True)
        if r.status_code == 200:
            return r.content
    except Exception:
        pass

    return None

# ── PPTX 생성 ─────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# 파스텔 배경 색상 리스트
BG_COLORS = [
    RGBColor(0xFF, 0xF0, 0xE6),  # 살구
    RGBColor(0xE6, 0xF5, 0xFF),  # 하늘
    RGBColor(0xE6, 0xFF, 0xED),  # 민트
    RGBColor(0xFF, 0xFF, 0xE6),  # 레몬
    RGBColor(0xF5, 0xE6, 0xFF),  # 라벤더
    RGBColor(0xFF, 0xE6, 0xF5),  # 핑크
    RGBColor(0xE6, 0xFF, 0xFF),  # 아쿠아
    RGBColor(0xFF, 0xFF, 0xE6),  # 크림
]

# 폰트 경로 (한국어 지원)
FONT_PATH = None
for fp in [
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    "/Library/Fonts/Arial Unicode MS.ttf",
]:
    if os.path.exists(fp):
        FONT_PATH = fp
        break

def make_pptx(title, slides_data, image_bytes_list):
    """슬라이드 데이터 + 이미지로 PPTX 생성 → 파일 경로 반환"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃

    # ── 표지 슬라이드 ─────────────────────────────────────────
    cover = prs.slides.add_slide(blank_layout)
    bg = cover.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xEE)

    # 제목 텍스트
    txBox = cover.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"📚 {title}"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)

    sub = tf.add_paragraph()
    sub.alignment = PP_ALIGN.CENTER
    subrun = sub.add_run()
    subrun.text = "함께 읽어요! 🌟"
    subrun.font.size = Pt(28)
    subrun.font.color.rgb = RGBColor(0xF0, 0x80, 0x40)

    # ── 콘텐츠 슬라이드 ──────────────────────────────────────
    for i, (slide_info, img_bytes) in enumerate(zip(slides_data, image_bytes_list)):
        slide = prs.slides.add_slide(blank_layout)

        # 배경색
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLORS[i % len(BG_COLORS)]

        # 이미지 (상단 75%)
        IMG_H = Inches(5.3)
        if img_bytes:
            try:
                img_io = BytesIO(img_bytes)
                # PIL로 크기 확인
                pil_img = Image.open(BytesIO(img_bytes))
                iw, ih = pil_img.size
                aspect = iw / ih

                # 가로세로 비율에 맞게 조정
                max_w = Inches(12.5)
                max_h = IMG_H
                if aspect > (max_w / max_h):
                    draw_w = max_w
                    draw_h = int(max_w / aspect)
                else:
                    draw_h = max_h
                    draw_w = int(max_h * aspect)

                left = int((SLIDE_W - draw_w) / 2)
                top = Inches(0.3)

                img_io = BytesIO(img_bytes)
                slide.shapes.add_picture(img_io, left, top, draw_w, draw_h)
            except Exception as e:
                print(f"이미지 삽입 오류: {e}")

        # 캡션 (하단)
        CAPTION_TOP = Inches(5.8)
        CAPTION_H = Inches(1.5)
        txBox = slide.shapes.add_textbox(
            Inches(0.5), CAPTION_TOP, Inches(12.33), CAPTION_H
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = slide_info["caption"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x2C, 0x2C)

        # 슬라이드 번호
        num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.alignment = PP_ALIGN.RIGHT
        num_run = num_p.add_run()
        num_run.text = f"{i+1}/{len(slides_data)}"
        num_run.font.size = Pt(14)
        num_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # 저장
    out_path = tempfile.mktemp(suffix=".pptx")
    prs.save(out_path)
    return out_path

# ── 메인 봇 루프 ──────────────────────────────────────────────

def process_photo(chat_id, file_id):
    """사진 수신 → 분석 → PPTX → 전송"""
    send_message(chat_id, "📸 책 사진 받았어요! 분석 중... (30초~1분 소요)")

    try:
        # 1. 사진 다운로드
        img_bytes = download_photo(file_id)
        send_message(chat_id, "🤖 Claude가 내용 읽는 중...")

        # 2. Claude로 분석
        data = analyze_book_image(img_bytes)
        title = data.get("title", "책 이야기")
        slides = data.get("slides", [])[:8]

        send_message(chat_id, f"📖 제목: {title}\n🎨 이미지 {len(slides)}장 찾는 중...")

        # 3. 이미지 검색
        image_bytes_list = []
        for i, s in enumerate(slides):
            kw = s.get("keyword", "children book illustration")
            print(f"  이미지 검색 [{i+1}]: {kw}")
            img = search_image(kw)
            image_bytes_list.append(img)
            time.sleep(0.5)

        # 4. PPTX 생성
        send_message(chat_id, "📊 파워포인트 만드는 중...")
        pptx_path = make_pptx(title, slides, image_bytes_list)

        # 5. 전송
        send_document(chat_id, pptx_path, caption=f"📚 {title} - {len(slides)}장 슬라이드")
        os.remove(pptx_path)
        send_message(chat_id, "✅ 완료! 다른 책 사진을 보내주세요 😊")

    except json.JSONDecodeError as e:
        send_message(chat_id, f"⚠️ 분석 오류: JSON 파싱 실패\n{e}")
    except Exception as e:
        send_message(chat_id, f"❌ 오류 발생: {str(e)[:200]}")
        print(f"오류: {e}")

def main():
    print("📚 책 PPTX 봇 시작!")
    offset = 0

    # 이전 offset 로드
    offset_file = os.path.expanduser("~/book_bot_offset.json")
    if os.path.exists(offset_file):
        with open(offset_file) as f:
            offset = json.load(f).get("offset", 0)
        print(f"이전 offset: {offset}")

    while True:
        try:
            res = tg_get("getUpdates", {
                "timeout": 30,
                "offset": offset,
                "allowed_updates": "message"
            })

            if not res.get("ok"):
                time.sleep(5)
                continue

            for update in res.get("result", []):
                offset = update["update_id"] + 1
                msg = update.get("message", {})
                chat_id = msg.get("chat", {}).get("id")

                if not chat_id:
                    continue

                # 사진 처리
                if "photo" in msg:
                    # 가장 큰 해상도 선택
                    photo = max(msg["photo"], key=lambda p: p["file_size"])
                    process_photo(chat_id, photo["file_id"])

                # 텍스트 명령
                elif "text" in msg:
                    text = msg["text"].strip()
                    if text.startswith("/start"):
                        send_message(chat_id,
                            "👋 안녕하세요! 책 PPTX 봇이에요.\n\n"
                            "📸 책 페이지 사진을 보내주시면\n"
                            "4살 수준으로 요약해서\n"
                            "예쁜 파워포인트로 만들어 드려요! 📚"
                        )

            # offset 저장
            with open(offset_file, "w") as f:
                json.dump({"offset": offset}, f)

        except requests.exceptions.ConnectionError:
            time.sleep(10)
        except Exception as e:
            print(f"폴링 오류: {e}")
            time.sleep(5)

if __name__ == "__main__":
    main()
